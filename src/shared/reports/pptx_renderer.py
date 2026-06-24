# ruff: noqa: E402
"""PPTX slide generators and generate_pptx entry point."""

from __future__ import annotations

import logging
from io import BytesIO
from types import SimpleNamespace
from typing import Any

logger = logging.getLogger(__name__)

from shared.reports.deck_model import DeckData, Section
from shared.reports.extraction import _extract_deck_data, fit_text

# ── Colour palette — modern navy/blue (matches deck template) ────────────────
_NAVY = "0C1425"
_NAVY_MID = "162036"
_BLUE = "3B7DFA"
_BLUE_LIGHT = "E8F0FE"
_GREEN = "22C55E"
_GREEN_DARK = "15803D"
_AMBER = "F59E0B"
_RED = "EF4444"
_TEXT = "1A2236"
_TEXT_MID = "475569"
_TEXT_LIGHT = "94A3B8"
_BORDER = "E2E8F0"
_SURFACE = "F1F5F9"
_WHITE = "FFFFFF"

# Badge colours
_BADGE_STRONG_BG = "DCFCE7"
_BADGE_STRONG_FG = "15803D"
_BADGE_BULLISH_BG = "DBEAFE"
_BADGE_BULLISH_FG = "1D4ED8"
_BADGE_EXPENSIVE_BG = "FEF3C7"
_BADGE_EXPENSIVE_FG = "92400E"
_BADGE_MODERATE_BG = "F1F5F9"
_BADGE_MODERATE_FG = "475569"

# Opportunity/Risk column colours
_OPP_BG = "F0FDF4"
_OPP_BORDER = "BBF7D0"
_RISK_BG = "FEF2F2"
_RISK_BORDER = "FECACA"

_SIGNAL_COLOURS = {
    "BUY": (_GREEN_DARK, _BADGE_STRONG_BG),
    "HOLD": (_BADGE_BULLISH_FG, _BADGE_BULLISH_BG),
    "SELL": (_RED, _RISK_BG),
}


# ── PPTX slide generators ──────────────────────────────────────────────────────


def _pptx_slide_title(deck: DeckData, h: SimpleNamespace) -> None:
    s = h.add_slide(dark=True)
    h.slide_label(s, "EQUITY RESEARCH REPORT", light=True)
    h.text(s, 0.83, 1.2, 10, 1.0, deck.company_name, 52, bold=True, color=_WHITE, font=h.FONT)
    exchange_label = (
        f"NYSE: {deck.ticker}" if not deck.exchange else f"{deck.exchange}: {deck.ticker}"
    )
    sector_label = f" · {deck.sector}" if deck.sector else ""
    date_label = f" · {deck.analysis_date}" if deck.analysis_date else ""
    h.text(
        s,
        0.83,
        2.1,
        10,
        0.4,
        f"{exchange_label}{sector_label}{date_label}",
        18,
        color=_TEXT_LIGHT,
        font=h.FONT,
    )
    kpi_y = 5.5
    h.text(s, 0.83, kpi_y, 2, 0.2, "Recommendation", 10, color=_TEXT_LIGHT, font=h.FONT, bold=True)
    rec_color = (
        _BLUE if deck.recommendation == "HOLD" else _GREEN if deck.recommendation == "BUY" else _RED
    )
    h.text(
        s,
        0.83,
        kpi_y + 0.25,
        2,
        0.45,
        deck.recommendation,
        28,
        bold=True,
        color=rec_color,
        font=h.MONO,
    )
    h.text(s, 3.3, kpi_y, 2, 0.2, "Confidence", 10, color=_TEXT_LIGHT, font=h.FONT, bold=True)
    h.text(
        s,
        3.3,
        kpi_y + 0.25,
        2,
        0.45,
        f"{deck.confidence:.0%}",
        28,
        bold=True,
        color=_WHITE,
        font=h.MONO,
    )
    if deck.scenarios.get("base"):
        h.text(
            s, 5.8, kpi_y, 2.5, 0.2, "Median Target", 10, color=_TEXT_LIGHT, font=h.FONT, bold=True
        )
        h.text(
            s,
            5.8,
            kpi_y + 0.25,
            2.5,
            0.45,
            deck.scenarios["base"],
            28,
            bold=True,
            color=_WHITE,
            font=h.MONO,
        )
    h.footer(s, dark=True)


def _pptx_slide_key_metrics(deck: DeckData, h: SimpleNamespace) -> None:
    if not deck.kpi_chips:
        return
    s = h.add_slide()
    h.slide_label(s, "KEY METRICS")
    h.slide_title(s, "Performance Snapshot")
    chips = deck.kpi_chips[:4]
    n = len(chips)
    chip_w = {1: 5.0, 2: 4.0, 3: 3.3, 4: 2.7}.get(n, 2.7)
    gap = 0.3
    total_w = n * chip_w + (n - 1) * gap
    start_x = (13.33 - total_w) / 2
    for i, chip in enumerate(chips):
        x = start_x + i * (chip_w + gap)
        h.kpi_chip(
            s,
            x,
            2.4,
            chip_w,
            chip["label"],
            chip["value"],
            chip.get("context", ""),
            chip.get("positive", True),
        )
    h.footer(s)


def _pptx_slide_thesis(deck: DeckData, h: SimpleNamespace) -> None:
    if not deck.executive_summary:
        return
    s = h.add_slide()
    h.slide_label(s, "INVESTMENT THESIS")
    h.slide_title(s, "Executive Summary")
    thesis_top = 2.4
    thesis_h = 4.2
    border = s.shapes.add_shape(
        1, h.Inches(0.83), h.Inches(thesis_top), h.Inches(0.04), h.Inches(thesis_h)
    )
    border.fill.solid()
    border.fill.fore_color.rgb = h.rgb(_BLUE)
    border.line.fill.background()
    h.rounded_rect(s, 0.87, thesis_top, 11.5, thesis_h, _SURFACE)
    fitted_text, fitted_size = fit_text(deck.executive_summary, 10.8, thesis_h - 0.6)
    h.text(
        s,
        1.2,
        thesis_top + 0.3,
        10.8,
        thesis_h - 0.6,
        fitted_text,
        fitted_size,
        color=_TEXT,
        font=h.FONT,
    )
    h.footer(s)


def _pptx_slide_financials(deck: DeckData, h: SimpleNamespace) -> None:
    if not deck.financials:
        return
    s = h.add_slide()
    h.slide_label(s, "FINANCIAL PERFORMANCE")
    h.slide_title(s, "Earnings & Growth")
    rows = deck.financials[:8]
    tbl_shape = s.shapes.add_table(
        len(rows) + 1,
        3,
        h.Inches(0.83),
        h.Inches(2.4),
        h.Inches(11.5),
        h.Inches(min(4.2, 0.55 + len(rows) * 0.52)),
    )
    tbl = tbl_shape.table
    tbl.columns[0].width = h.Inches(5)
    tbl.columns[1].width = h.Inches(3.25)
    tbl.columns[2].width = h.Inches(3.25)
    for ci, hdr in enumerate(["METRIC", "CURRENT", "CONTEXT"]):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = h.rgb(_WHITE)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = h.Pt(10)
        run.font.color.rgb = h.rgb(_TEXT_LIGHT)
        run.font.name = h.FONT
        cell.margin_top = h.Pt(8)
        cell.margin_bottom = h.Pt(8)
        cell.margin_left = h.Pt(16)
    for ri, (metric, current, context) in enumerate(rows):
        for ci, (txt, is_metric) in enumerate([(metric, True), (current, False), (context, False)]):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = h.rgb(_WHITE)
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = txt
            run.font.size = h.Pt(14)
            run.font.name = h.MONO if not is_metric else h.FONT
            run.font.bold = is_metric
            run.font.color.rgb = h.rgb(
                _TEXT if is_metric else _GREEN_DARK if txt.startswith("+") else _TEXT
            )
            cell.margin_top = h.Pt(12)
            cell.margin_bottom = h.Pt(12)
            cell.margin_left = h.Pt(16)
    h.footer(s)


def _pptx_slide_valuation(deck: DeckData, h: SimpleNamespace) -> None:
    if not deck.valuation_table and not deck.scenarios:
        return
    s = h.add_slide()
    h.slide_label(s, "VALUATION")
    h.slide_title(s, "Scenario Analysis")

    # P5 dedup: determine scenario card labels from the scenarios dict
    _scenario_card_labels = set()
    scenario_list = []
    if deck.scenarios.get("bull"):
        scenario_list.append(("Bull Case", deck.scenarios["bull"], _GREEN_DARK))
        _scenario_card_labels.add("Bull Case Target")
    if deck.scenarios.get("base"):
        scenario_list.append(("Base Case", deck.scenarios["base"], _BLUE))
        _scenario_card_labels.add("Base Case")
    if deck.scenarios.get("dcf"):
        scenario_list.append(("DCF Fair Value", deck.scenarios["dcf"], _AMBER))
        _scenario_card_labels.add("DCF Fair Value")
    if deck.scenarios.get("bear"):
        scenario_list.append(("Bear Case", deck.scenarios["bear"], _RED))
        _scenario_card_labels.add("Bear Case Target")

    # Cards preferred when ≥2 scenarios; skip those rows from valuation table
    show_cards = len(scenario_list) >= 2
    if deck.valuation_table:
        vtbl_rows = [
            r
            for r in deck.valuation_table[:10]
            if not show_cards or r[0] not in _scenario_card_labels
        ][:6]
        vtbl_shape = s.shapes.add_table(
            len(vtbl_rows) + 1,
            2,
            h.Inches(0.83),
            h.Inches(2.4),
            h.Inches(5.5),
            h.Inches(min(4.0, 0.55 + len(vtbl_rows) * 0.52)),
        )
        vtbl = vtbl_shape.table
        vtbl.columns[0].width = h.Inches(3)
        vtbl.columns[1].width = h.Inches(2.5)
        for ci, hdr in enumerate(["METRIC", "VALUE"]):
            cell = vtbl.cell(0, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = h.rgb(_WHITE)
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = hdr
            run.font.bold = True
            run.font.size = h.Pt(10)
            run.font.color.rgb = h.rgb(_TEXT_LIGHT)
            run.font.name = h.FONT
            cell.margin_left = h.Pt(16)
        for ri, (metric, value) in enumerate(vtbl_rows):
            for ci, (txt, is_metric) in enumerate([(metric, True), (value, False)]):
                cell = vtbl.cell(ri + 1, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = h.rgb(_WHITE)
                p = cell.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = txt
                run.font.size = h.Pt(13)
                run.font.name = h.MONO if not is_metric else h.FONT
                run.font.bold = is_metric
                val_color = _GREEN_DARK if txt.startswith("+") else _TEXT
                run.font.color.rgb = h.rgb(_TEXT if is_metric else val_color)
                cell.margin_top = h.Pt(10)
                cell.margin_bottom = h.Pt(10)
                cell.margin_left = h.Pt(16)
    if show_cards:
        card_x = 7.0
        card_w = 5.5
        card_h = 0.95
        for i, (label, value, color) in enumerate(scenario_list):
            cy = 2.4 + i * (card_h + 0.2)
            h.rounded_rect(s, card_x, cy, card_w, card_h, _SURFACE)
            h.text(
                s,
                card_x + 0.3,
                cy + 0.12,
                card_w - 0.6,
                0.2,
                label.upper(),
                10,
                color=_TEXT_LIGHT,
                font=h.FONT,
                bold=True,
            )
            h.text(
                s,
                card_x + 0.3,
                cy + 0.38,
                card_w - 0.6,
                0.45,
                value,
                28,
                bold=True,
                color=color,
                font=h.MONO,
            )
    h.footer(s)


def _pptx_slide_scorecard(deck: DeckData, h: SimpleNamespace) -> None:
    if len(deck.scorecard) <= 1:
        return
    s = h.add_slide()
    h.slide_label(s, "ASSESSMENT")
    h.slide_title(s, "Investment Scorecard")
    badge_colors = {
        "strong": (_BADGE_STRONG_FG, _BADGE_STRONG_BG),
        "bullish": (_BADGE_BULLISH_FG, _BADGE_BULLISH_BG),
        "expensive": (_BADGE_EXPENSIVE_FG, _BADGE_EXPENSIVE_BG),
        "moderate": (_BADGE_MODERATE_FG, _BADGE_MODERATE_BG),
    }
    items = deck.scorecard[:6]
    n_items = len(items)
    if n_items <= 3:
        cols = n_items
    elif n_items == 4:
        cols = 2
    else:
        cols = 3
    card_w = 3.5
    card_h = 1.2
    gap_x = 0.35
    gap_y = 0.25
    total_w = cols * card_w + (cols - 1) * gap_x
    start_x = (13.33 - total_w) / 2
    start_y = 2.4
    for i, (dim, rating, badge_type) in enumerate(items):
        col = i % cols
        row = i // cols
        cx = start_x + col * (card_w + gap_x)
        cy = start_y + row * (card_h + gap_y)
        h.rounded_rect(s, cx, cy, card_w, card_h, _SURFACE)
        h.text(
            s,
            cx + 0.22,
            cy + 0.18,
            card_w - 0.4,
            0.2,
            dim,
            11,
            color=_TEXT_MID,
            font=h.FONT,
            bold=True,
        )
        fg, bg = badge_colors.get(badge_type, badge_colors["moderate"])
        h.rounded_rect(s, cx + 0.22, cy + 0.55, 1.4, 0.38, bg)
        h.text(
            s,
            cx + 0.22,
            cy + 0.55,
            1.4,
            0.38,
            rating,
            12,
            bold=True,
            color=fg,
            font=h.FONT,
            align=h.PP_ALIGN.CENTER,
            anchor=h.MSO_ANCHOR.MIDDLE,
        )
    h.footer(s)


def _pptx_slide_peers(deck: DeckData, h: SimpleNamespace) -> None:
    if not deck.peer_names:
        return
    s = h.add_slide()
    h.slide_label(s, "COMPETITIVE LANDSCAPE")
    h.slide_title(s, "Peer Comparison")
    if not deck.peers:
        peer_text = f"Identified peers: {', '.join(deck.peer_names[:5])}"
        h.rounded_rect(s, 0.83, 2.8, 11.5, 1.5, _SURFACE)
        h.text(s, 1.2, 3.0, 10.8, 1.0, peer_text, 18, color=_TEXT, font=h.FONT)
        h.text(
            s,
            1.2,
            3.8,
            10.8,
            0.5,
            "Detailed metrics not available for this analysis.",
            14,
            color=_TEXT_LIGHT,
            font=h.FONT,
        )
        h.footer(s)
        return
    headers = ["METRIC", deck.ticker] + deck.peer_names[:2]
    n_cols = len(headers)
    ptbl_shape = s.shapes.add_table(
        len(deck.peers) + 1,
        n_cols,
        h.Inches(0.83),
        h.Inches(2.4),
        h.Inches(11.5),
        h.Inches(min(4.2, 0.55 + len(deck.peers) * 0.52)),
    )
    ptbl = ptbl_shape.table
    col_w = 11.5 / n_cols
    for ci in range(n_cols):
        ptbl.columns[ci].width = h.Inches(col_w if ci > 0 else col_w * 1.3)
    for ci, hdr in enumerate(headers):
        cell = ptbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = h.rgb(_WHITE)
        p = cell.text_frame.paragraphs[0]
        p.alignment = h.PP_ALIGN.CENTER if ci > 0 else h.PP_ALIGN.LEFT
        run = p.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = h.Pt(10)
        run.font.color.rgb = h.rgb(_TEXT_LIGHT)
        run.font.name = h.FONT
        cell.margin_left = h.Pt(16)
    for ri, peer_row in enumerate(deck.peers[:6]):
        values = [peer_row.get("metric", "")] + [
            peer_row.get(f"col{ci}", "N/A") for ci in range(n_cols - 1)
        ]
        is_highlight = ri % 2 == 0
        for ci, txt in enumerate(values):
            cell = ptbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = h.rgb(_BLUE_LIGHT if is_highlight else _WHITE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = h.PP_ALIGN.CENTER if ci > 0 else h.PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(txt)
            run.font.size = h.Pt(13)
            run.font.name = h.MONO if ci > 0 else h.FONT
            run.font.bold = ci == 0 or (ci == 1 and is_highlight)
            run.font.color.rgb = h.rgb(_TEXT if txt != "N/A" else _TEXT_LIGHT)
            cell.margin_top = h.Pt(10)
            cell.margin_bottom = h.Pt(10)
            cell.margin_left = h.Pt(16)
    h.footer(s)


def _pptx_slide_risk_reward(deck: DeckData, h: SimpleNamespace) -> None:
    if not deck.opportunities and not deck.risks:
        return
    s = h.add_slide()
    h.slide_label(s, "RISK–REWARD")
    h.slide_title(s, "Assessment")
    col_w = 5.5
    gap = 0.5
    total = 2 * col_w + gap
    sx = (13.33 - total) / 2
    n_items = max(len(deck.opportunities[:5]), len(deck.risks[:5]), 1)
    item_h = 0.7
    col_h = 0.9 + n_items * item_h
    opp_shape = s.shapes.add_shape(
        h.MSO_SHAPE.ROUNDED_RECTANGLE,
        h.Inches(sx),
        h.Inches(2.4),
        h.Inches(col_w),
        h.Inches(col_h),
    )
    opp_shape.fill.solid()
    opp_shape.fill.fore_color.rgb = h.rgb(_OPP_BG)
    opp_shape.line.color.rgb = h.rgb(_OPP_BORDER)
    opp_shape.line.width = h.Pt(1.5)
    h.text(
        s,
        sx + 0.3,
        2.6,
        col_w - 0.6,
        0.35,
        "Growth Opportunities",
        16,
        bold=True,
        color=_GREEN_DARK,
        font=h.FONT,
    )
    for i, opp in enumerate(deck.opportunities[:5]):
        opp_text = (opp[:117] + "...") if len(opp) > 120 else opp
        bullet_y = 3.15 + i * item_h
        dot = s.shapes.add_shape(
            h.MSO_SHAPE.OVAL,
            h.Inches(sx + 0.4),
            h.Inches(bullet_y + 0.06),
            h.Inches(0.08),
            h.Inches(0.08),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = h.rgb(_GREEN)
        dot.line.fill.background()
        h.text(s, sx + 0.65, bullet_y, col_w - 1.0, item_h, opp_text, 12, color=_TEXT, font=h.FONT)
    rx = sx + col_w + gap
    risk_shape = s.shapes.add_shape(
        h.MSO_SHAPE.ROUNDED_RECTANGLE,
        h.Inches(rx),
        h.Inches(2.4),
        h.Inches(col_w),
        h.Inches(col_h),
    )
    risk_shape.fill.solid()
    risk_shape.fill.fore_color.rgb = h.rgb(_RISK_BG)
    risk_shape.line.color.rgb = h.rgb(_RISK_BORDER)
    risk_shape.line.width = h.Pt(1.5)
    h.text(s, rx + 0.3, 2.6, col_w - 0.6, 0.35, "Key Risks", 16, bold=True, color=_RED, font=h.FONT)
    for i, risk in enumerate(deck.risks[:5]):
        risk_text = (risk[:117] + "...") if len(risk) > 120 else risk
        bullet_y = 3.15 + i * item_h
        dot = s.shapes.add_shape(
            h.MSO_SHAPE.OVAL,
            h.Inches(rx + 0.4),
            h.Inches(bullet_y + 0.06),
            h.Inches(0.08),
            h.Inches(0.08),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = h.rgb(_RED)
        dot.line.fill.background()
        h.text(s, rx + 0.65, bullet_y, col_w - 1.0, item_h, risk_text, 12, color=_TEXT, font=h.FONT)
    h.footer(s)


def _pptx_slide_extra(deck: DeckData, h: SimpleNamespace, section: Section) -> None:
    from shared.reports.deck_model import _MAX_SLIDE_CHARS

    if not section.body:
        return
    body = section.body
    chunks: list[str] = []
    while len(body) > _MAX_SLIDE_CHARS:
        split_at = body.rfind("\n", 0, _MAX_SLIDE_CHARS)
        if split_at < 0:
            split_at = _MAX_SLIDE_CHARS
        chunks.append(body[:split_at])
        body = body[split_at:].strip()
    chunks.append(body)
    for ci, chunk in enumerate(chunks):
        sld = h.add_slide()
        label_text = section.title.upper()
        h.slide_label(sld, label_text if ci == 0 else f"{label_text} (CONT.)")
        h.slide_title(sld, section.title if ci == 0 else f"{section.title} (cont.)")
        border = sld.shapes.add_shape(
            1, h.Inches(0.83), h.Inches(2.5), h.Inches(0.04), h.Inches(4.0)
        )
        border.fill.solid()
        border.fill.fore_color.rgb = h.rgb(_BLUE)
        border.line.fill.background()
        h.rounded_rect(sld, 0.87, 2.5, 11.5, 4.0, _SURFACE)
        h.text(sld, 1.2, 2.7, 10.8, 3.6, chunk, 14, color=_TEXT, font=h.FONT)
        h.footer(sld)


def _pptx_slide_conclusion(deck: DeckData, h: SimpleNamespace) -> None:
    s = h.add_slide(dark=True)
    h.slide_label(s, "CONCLUSION", left=0.83, top=2.0, light=True)
    div = s.shapes.add_shape(1, h.Inches(6.2), h.Inches(2.5), h.Inches(0.9), h.Inches(0.04))
    div.fill.solid()
    div.fill.fore_color.rgb = h.rgb(_BLUE)
    div.line.fill.background()
    h.text(
        s,
        0,
        2.8,
        13.33,
        0.7,
        deck.recommendation,
        44,
        bold=True,
        color=_BLUE,
        font=h.MONO,
        align=h.PP_ALIGN.CENTER,
    )
    h.text(
        s,
        0,
        3.5,
        13.33,
        0.4,
        f"{deck.confidence:.0%} Confidence",
        14,
        color=_TEXT_LIGHT,
        font=h.FONT,
        align=h.PP_ALIGN.CENTER,
    )
    conclusion_text = deck.executive_summary[:300] if deck.executive_summary else ""
    if conclusion_text:
        h.text(
            s,
            2.5,
            4.2,
            8.33,
            1.5,
            conclusion_text,
            15,
            color=_TEXT_LIGHT,
            font=h.FONT,
            align=h.PP_ALIGN.CENTER,
        )
    h.footer(s, dark=True)


# ── PPTX generation ───────────────────────────────────────────────────────────


def generate_pptx(
    brief_data: dict[str, Any],
    ticker: str,
    recommendation: str,
    confidence: float,
    analysis_date: str,
    company_info: dict[str, Any] | None = None,
) -> BytesIO:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    def rgb(hex_str: str) -> RGBColor:
        return RGBColor.from_string(hex_str)  # type: ignore[no-any-return,no-untyped-call]

    logger.info(
        "Generating PPTX report for %s (recommendation=%s, confidence=%.0f%%)",
        ticker,
        recommendation,
        confidence * 100,
    )
    deck = _extract_deck_data(
        brief_data, ticker, recommendation, confidence, analysis_date, company_info=company_info
    )
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    FONT = "Calibri"
    MONO = "Consolas"

    def add_slide(dark: bool = False) -> Any:
        sld = prs.slides.add_slide(blank_layout)
        bg = sld.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = rgb(_NAVY if dark else _WHITE)
        bg.line.fill.background()
        return sld

    def _text(
        slide: Any,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        size: int,
        bold: bool = False,
        color: str = _TEXT,
        align: Any = PP_ALIGN.LEFT,
        font: str = FONT,
        wrap: bool = True,
        anchor: Any = MSO_ANCHOR.TOP,
    ) -> Any:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = wrap
        tf.auto_size = None
        try:
            tf.vertical_anchor = anchor
        except Exception:
            logger.debug("Could not set vertical_anchor on text frame")
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        run.font.name = font
        return box

    def _multiline(
        slide: Any, left: float, top: float, width: float, height: float,
        lines: list[str], size: int, color: str = _TEXT, font: str = FONT,
        line_spacing: float = 1.5
    ) -> Any:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line_text in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(4)
            run = p.add_run()
            run.text = line_text
            run.font.size = Pt(size)
            run.font.color.rgb = rgb(color)
            run.font.name = font
        return box

    def _rounded_rect(
        slide: Any, left: float, top: float, width: float, height: float,
        fill_color: str, border_color: str | None = None
    ) -> Any:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_color)
        if border_color:
            shape.line.color.rgb = rgb(border_color)
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    def _slide_label(
        slide: Any, text: str, left: float = 0.83, top: float = 0.7,
        light: bool = False
    ) -> None:
        _text(
            slide,
            left,
            top,
            4,
            0.25,
            text,
            11,
            bold=True,
            color=_TEXT_LIGHT if light else _BLUE,
            font=FONT,
        )

    def _slide_title(
        slide: Any, text: str, left: float = 0.83, top: float = 0.95,
        light: bool = False
    ) -> None:
        _text(
            slide,
            left,
            top,
            10,
            0.7,
            text,
            36,
            bold=True,
            color=_WHITE if light else _TEXT,
            font=FONT,
        )

    def _footer(slide: Any, dark: bool = False) -> None:
        color = _TEXT_LIGHT if not dark else _TEXT_MID
        _text(slide, 0.83, 7.0, 5, 0.3, "© 2026 Institutional Equity Research", 9, color=color)
        _text(
            slide,
            7,
            7.0,
            5.5,
            0.3,
            "For informational purposes only",
            9,
            color=color,
            align=PP_ALIGN.RIGHT,
        )

    def _kpi_chip(
        slide: Any, left: float, top: float, width: float, label: str,
        value: str, context: str = "", positive: bool = True, dark: bool = False
    ) -> None:
        chip_bg = _NAVY_MID if dark else _SURFACE
        _rounded_rect(slide, left, top, width, 1.2, chip_bg)
        _text(
            slide,
            left + 0.22,
            top + 0.2,
            width - 0.4,
            0.25,
            label,
            11,
            color=_TEXT_MID if not dark else _TEXT_LIGHT,
            font=FONT,
            bold=True,
        )
        val_color = _GREEN_DARK if positive else _RED
        if not positive and value.startswith("-"):
            val_color = _RED
        _text(
            slide,
            left + 0.22,
            top + 0.5,
            width - 0.4,
            0.45,
            value,
            30,
            bold=True,
            color=val_color,
            font=MONO,
        )
        if context:
            _text(slide, left + 0.22, top + 0.92, width - 0.4, 0.2, context, 10, color=_TEXT_LIGHT)

    h = SimpleNamespace(
        add_slide=add_slide,
        text=_text,
        multiline=_multiline,
        rounded_rect=_rounded_rect,
        kpi_chip=_kpi_chip,
        slide_label=_slide_label,
        slide_title=_slide_title,
        footer=_footer,
        rgb=rgb,
        FONT=FONT,
        MONO=MONO,
        Inches=Inches,
        Pt=Pt,
        Emu=Emu,
        PP_ALIGN=PP_ALIGN,
        MSO_ANCHOR=MSO_ANCHOR,
        MSO_SHAPE=MSO_SHAPE,
        blank_layout=blank_layout,
        prs=prs,
    )

    _pptx_slide_title(deck, h)
    logger.debug("PPTX slide: title")
    _pptx_slide_key_metrics(deck, h)
    logger.debug("PPTX slide: key_metrics")
    _pptx_slide_thesis(deck, h)
    logger.debug("PPTX slide: thesis")
    _pptx_slide_financials(deck, h)
    logger.debug("PPTX slide: financials")
    _pptx_slide_valuation(deck, h)
    logger.debug("PPTX slide: valuation")
    _pptx_slide_scorecard(deck, h)
    logger.debug("PPTX slide: scorecard")
    _pptx_slide_peers(deck, h)
    logger.debug("PPTX slide: peers")
    _pptx_slide_risk_reward(deck, h)
    logger.debug("PPTX slide: risk_reward")
    for section in deck.sections[:4]:
        _pptx_slide_extra(deck, h, section)
        logger.debug("PPTX slide: extra (%s)", section.title)
    _pptx_slide_conclusion(deck, h)
    logger.debug("PPTX slide: conclusion")

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    logger.info("PPTX report generated for %s (%d bytes)", ticker, buf.getbuffer().nbytes)
    return buf
