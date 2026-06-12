# ruff: noqa: E402
"""Playwright-based HTML → PPTX/PDF export using deck-stage's built-in APIs.

On Windows the server's SelectorEventLoop cannot spawn subprocesses, so
Playwright is launched in a child process that uses ProactorEventLoop.
"""

from __future__ import annotations

import asyncio
import logging
import sys
from concurrent.futures import ProcessPoolExecutor
from io import BytesIO

logger = logging.getLogger(__name__)

_executor = ProcessPoolExecutor(max_workers=2)


def _set_proactor_policy() -> None:
    """Ensure the child process uses ProactorEventLoop on Windows."""
    if sys.platform == "win32":
        asyncio.set_event_loop_policy(asyncio.WindowsProactorEventLoopPolicy())


def _pptx_worker(html: str) -> bytes:
    """Run in a child process — has its own ProactorEventLoop on Windows."""
    _set_proactor_policy()
    from io import BytesIO as _BytesIO

    from playwright.sync_api import sync_playwright

    html = html.replace("<deck-stage ", "<deck-stage noscale ")

    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": 1920, "height": 1080})
        page.set_content(html, wait_until="networkidle")
        page.wait_for_function(
            """() => {
                const ds = document.querySelector('deck-stage');
                if (!ds || !customElements.get('deck-stage')) return false;
                if (ds.hasAttribute('data-fonts-pending')) return false;
                return ds._slides && ds._slides.length > 0;
            }""",
            timeout=15000,
        )
        total = page.evaluate("document.querySelector('deck-stage')._slides.length")

        from pptx import Presentation
        from pptx.util import Emu

        prs = Presentation()
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)
        blank_layout = prs.slide_layouts[6]

        for i in range(total):
            page.evaluate(f"document.querySelector('deck-stage').goTo({i})")
            page.wait_for_timeout(150)
            section = page.query_selector("section[data-deck-active]")
            if section:
                screenshot = section.screenshot(type="png")
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    _BytesIO(screenshot), Emu(0), Emu(0), Emu(12192000), Emu(6858000)
                )

        browser.close()

        buf = _BytesIO()
        prs.save(buf)
        return buf.getvalue()


def _pdf_worker(html: str) -> bytes:
    """Run in a child process — has its own ProactorEventLoop on Windows."""
    _set_proactor_policy()
    from playwright.sync_api import sync_playwright

    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": 1920, "height": 1080})
        page.set_content(html, wait_until="networkidle")
        page.emulate_media(media="print")
        pdf_bytes = page.pdf(
            landscape=True, width="1920px", height="1080px", print_background=True
        )
        browser.close()
        return pdf_bytes


async def html_to_pptx(html: str) -> BytesIO:
    """Convert HTML deck to PPTX via headless Chromium screenshots."""
    loop = asyncio.get_running_loop()
    data = await loop.run_in_executor(_executor, _pptx_worker, html)
    return BytesIO(data)


async def html_to_pdf(html: str) -> BytesIO:
    """Convert HTML deck to PDF via headless Chromium print mode."""
    loop = asyncio.get_running_loop()
    data = await loop.run_in_executor(_executor, _pdf_worker, html)
    return BytesIO(data)


def html_to_pptx_sync(html: str) -> BytesIO:
    return BytesIO(_pptx_worker(html))


def html_to_pdf_sync(html: str) -> BytesIO:
    return BytesIO(_pdf_worker(html))
